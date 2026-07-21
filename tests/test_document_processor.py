"""Smoke tests for document_processor — one per supported format.

Each test creates a tiny in-memory file using the same library the
processor uses for extraction, writes it to a temp path, and verifies
that ``extract_text`` returns the expected plain text.
"""

import tempfile
from pathlib import Path

import pytest

from app.ingestion.document_processor import extract_text, is_supported


# ── is_supported ──────────────────────────────────────────────────────────


class TestIsSupported:
    def test_pdf(self):
        assert is_supported("report.pdf") is True

    def test_docx(self):
        assert is_supported("letter.docx") is True

    def test_xlsx(self):
        assert is_supported("data.xlsx") is True

    def test_xls(self):
        assert is_supported("legacy.xls") is True

    def test_pptx(self):
        assert is_supported("slides.pptx") is True

    def test_txt_rejected(self):
        assert is_supported("notes.txt") is False

    def test_no_extension(self):
        assert is_supported("Makefile") is False


# ── PDF extraction ────────────────────────────────────────────────────────


class TestExtractPdf:
    def test_simple_pdf(self):
        """Write a minimal valid PDF as raw bytes and extract its text."""
        # A minimal PDF with content stream containing "Hello World"
        content = b"BT /F1 12 Tf 100 700 Td (Hello from pypdf) Tj ET"
        pdf_bytes = (
            b"%PDF-1.4\n"
            b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
            b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
            b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]"
            b" /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
            b"4 0 obj\n<< /Length " + str(len(content)).encode() + b" >>\nstream\n"
            + content + b"\nendstream\nendobj\n"
            b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
            b"xref\n0 6\n"
            b"0000000000 65535 f \n"
            b"0000000009 00000 n \n"
            b"0000000058 00000 n \n"
            b"0000000115 00000 n \n"
            b"0000000266 00000 n \n"
            b"0000000341 00000 n \n"
            b"trailer\n<< /Size 6 /Root 1 0 R >>\n"
            b"startxref\n393\n"
            b"%%EOF"
        )

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            path = f.name
            f.write(pdf_bytes)

        try:
            text = extract_text(path, "test.pdf")
            assert "Hello from pypdf" in text
        finally:
            Path(path).unlink(missing_ok=True)


# ── DOCX extraction ───────────────────────────────────────────────────────


class TestExtractDocx:
    def test_paragraphs(self, tmp_dir: Path):
        """Create a .docx with one paragraph and extract its text."""
        from docx import Document

        path = tmp_dir / "test.docx"
        doc = Document()
        doc.add_paragraph("Hello from python-docx")
        doc.save(str(path))

        text = extract_text(str(path), "document.docx")
        assert "Hello from python-docx" in text

    def test_table(self, tmp_dir: Path):
        """Create a .docx with a simple table."""
        from docx import Document

        path = tmp_dir / "table.docx"
        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"
        doc.save(str(path))

        text = extract_text(str(path), "table.docx")
        assert "A1 | B1" in text
        assert "A2 | B2" in text


# ── XLSX extraction ───────────────────────────────────────────────────────


class TestExtractXlsx:
    def test_simple_sheet(self, tmp_dir: Path):
        """Create a single-sheet workbook and extract its text."""
        import openpyxl

        path = tmp_dir / "test.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Name"
        ws["B1"] = "Age"
        ws["A2"] = "Alice"
        ws["B2"] = "30"
        wb.save(str(path))

        text = extract_text(str(path), "data.xlsx")
        assert "Sheet: Data" in text
        assert "Name | Age" in text
        assert "Alice | 30" in text

    def test_multiple_sheets(self, tmp_dir: Path):
        import openpyxl

        path = tmp_dir / "multi.xlsx"
        wb = openpyxl.Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1["A1"] = "One"
        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = "Two"
        wb.save(str(path))

        text = extract_text(str(path), "multi.xlsx")
        assert "Sheet: Sheet1" in text
        assert "One" in text
        assert "Sheet: Sheet2" in text
        assert "Two" in text


# ── PPTX extraction ───────────────────────────────────────────────────────


class TestExtractPptx:
    def test_single_slide(self, tmp_dir: Path):
        """Create a one-slide presentation and extract its text."""
        from pptx import Presentation
        from pptx.util import Inches

        path = tmp_dir / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = tx_box.text_frame
        tf.text = "Hello from python-pptx"
        prs.save(str(path))

        text = extract_text(str(path), "slides.pptx")
        assert "Hello from python-pptx" in text
