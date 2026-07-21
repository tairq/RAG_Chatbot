"""
Multi-format document text extraction.

Supports:
  - PDF  (.pdf)   — via pypdf
  - DOCX (.docx)  — via python-docx
  - XLSX (.xlsx)  — via openpyxl
  - PPTX (.pptx)  — via python-pptx (optional, graceful fallback)

Each extractor returns a plain-text string suitable for the chunker.
"""

import logging
import zipfile
from pathlib import Path
from xml.etree import ElementTree

logger = logging.getLogger("rag_chatbot.document_processor")

# ── Supported extensions ──────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xls", ".pptx"}


# ── Extension helpers ────────────────────────────────────────────────


def _ext(filename: str) -> str:
    """Return the lower-case extension of *filename*."""
    return Path(filename).suffix.lower()


# ── PDF extraction (unchanged from pdf_processor.py) ─────────────────


def _extract_pdf(file_path: str) -> str:
    """Extract text from a PDF via pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    parts: list[str] = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            parts.append(page_text)
    return "\n".join(parts)


# ── DOCX extraction ──────────────────────────────────────────────────


def _extract_docx(file_path: str) -> str:
    """Extract text from a .docx file via python-docx."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx is required to process .docx files. "
            "Install it with: pip install python-docx"
        )

    doc = Document(file_path)
    parts: list[str] = []

    # Paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # Tables (row by row, cell by cell)
    for table in doc.tables:
        rows_text: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows_text.append(" | ".join(cells))
        parts.append("\n".join(rows_text))

    return "\n\n".join(parts)


# ── XLSX / XLS extraction ────────────────────────────────────────────


def _extract_xlsx(file_path: str) -> str:
    """Extract text from an Excel workbook via openpyxl.

    Each sheet is rendered as a text block with rows joined by newlines
    and cells separated by pipe characters.  Empty rows are skipped.
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError(
            "openpyxl is required to process .xlsx/.xls files. "
            "Install it with: pip install openpyxl"
        )

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_parts: list[str] = [f"--- Sheet: {sheet_name} ---"]

        row_count = 0
        for row in ws.iter_rows(values_only=True):
            # Skip fully empty rows
            if all(cell is None or (isinstance(cell, str) and cell.strip() == "") for cell in row):
                continue
            cleaned = [str(cell).strip() if cell is not None else "" for cell in row]
            sheet_parts.append(" | ".join(cleaned))
            row_count += 1

        if row_count > 0:
            parts.append("\n".join(sheet_parts))

    wb.close()
    return "\n\n".join(parts)


# ── PPTX extraction ──────────────────────────────────────────────────


def _extract_pptx(file_path: str) -> str:
    """Extract text from a .pptx file.

    Uses python-pptx if installed; otherwise falls back to a basic
    XML-parsing approach that extracts text from slide shapes.
    """
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        parts: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_parts.append(" | ".join(cells))
            parts.append("\n".join(slide_parts))

        return "\n\n".join(parts)

    except ImportError:
        logger.info("python-pptx not available; using XML fallback for .pptx")
        return _extract_pptx_xml_fallback(file_path)


def _extract_pptx_xml_fallback(file_path: str) -> str:
    """Basic PPTX text extraction via XML parsing (no python-pptx)."""
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    parts: list[str] = []

    try:
        with zipfile.ZipFile(file_path) as z:
            # Discover slides
            slide_paths = sorted(
                [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            )

            for slide_path in slide_paths:
                slide_num = slide_path.replace("ppt/slides/slide", "").replace(".xml", "")
                slide_parts: list[str] = [f"--- Slide {slide_num} ---"]

                xml_content = z.read(slide_path)
                root = ElementTree.fromstring(xml_content)

                for t_elem in root.iter(f"{{{ns['a']}}}t"):
                    text = (t_elem.text or "").strip()
                    if text:
                        slide_parts.append(text)

                parts.append("\n".join(slide_parts))

        return "\n\n".join(parts)

    except Exception as exc:
        logger.warning("XML fallback for .pptx failed: %s", exc)
        return ""


# ── Public API ───────────────────────────────────────────────────────


def extract_text(file_path: str, filename: str) -> str:
    """Extract text from a document based on its file extension.

    Args:
        file_path: Path to the saved file on disk.
        filename:  Original filename (used to detect extension).

    Returns:
        Extracted plain text.

    Raises:
        ValueError: If the file extension is not supported.
    """
    ext = _ext(filename)

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext in (".xlsx", ".xls"):
        return _extract_xlsx(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    else:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )


def is_supported(filename: str) -> bool:
    """Return ``True`` if *filename* has a supported extension."""
    return _ext(filename) in SUPPORTED_EXTENSIONS
